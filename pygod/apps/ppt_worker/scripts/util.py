import os
import glob
import json
from pptx import Presentation
import win32com.client

def get_album_tmp(script):
    try:
        album = script.rsplit('\n')[1]
        if len(album)>15:
            return 'NA'
        else:
            return album
    except:
        return 'NA'

def get_group_id(script, song_id):
    try:
        album = script.rsplit('\n')[1]
        if len(album)>15:
            album = 'NA' 
    except:
        album = 'NA'
    return f'{song_id}_{album}_NA'


#key: collection to be transferred to, value: collection being transferred from
mapping_rule_temp = {'name':'song_id',
                'script_text_for_ppt_worker':'script',
                'album':"get_album_tmp+script",
                'band':'NA',
                'script_text_link':'nan',
                'script_piano_link':'nan',
                'mp3_link_dropbox':"nan",
                'vedio_link_youtube':"nan",
                'archive_date':"2023-11-14",
                'archive_worker':"灿荣",
                'notes':"nan",
                'checked':False,
                'corrected':False,
                'group_id':'get_group_id+script+song_id'
                }
def transfer_collection_docs_to_db(mongo_client, db_name_from='ccg-PPT', collection_name_from='song_info', db_name_to='ccg-hymn', collection_name_to='hymn_info', mapping_rule=mapping_rule_temp, limits = None):
    collection_data_from = list(mongo_client[db_name_from][collection_name_from].find({}))
    if limits!=None:
        assert type(limits)==int, "limits has to be of type int"
        collection_data_from = collection_data_from[0:limits]
    collection_data_to = []
    for each in collection_data_from:
        doc = {}
        for value, key in mapping_rule.items():
            if key in each:
                doc[value] = each[key]
            elif key=='nan':
                doc[value] = 'na'
            elif type(key)==str and key.startswith('get_'):
                func, *args = key.rsplit('+')
                args_decoded = []
                for arg in args:
                    if arg in each:
                        args_decoded.append(each[arg])
                    else:
                        args_decoded.append(arg)
                doc[value] = eval(func)(*args_decoded)
            else:
                doc[value] = key
        collection_data_to.append(doc)
    mongo_client[db_name_to][collection_name_to].insert_many(collection_data_to)
    return collection_data_to

def extract_text_from_ppt_file(ppt_file_path):
    assert os.path.exists(ppt_file_path), "The specified ppt file path is not existing!"
    _, name = os.path.split(ppt_file_path)
    name = name.rsplit('.')[0]
    #assert os.path.exists(txt_file_folder), "the specified txt file folder is not existing!"
    ppt = Presentation(ppt_file_path)
    txt_box = [name]
    for slide in list(ppt.slides)[1:]:
        txt_temp = []
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.text.startswith('启应经文'):
                txt_temp.append(shape.text.replace('\u3000','　'))
        if len(txt_temp)==3:
            txt_box.append(''.join(txt_temp[0:2]))
            txt_box.append(txt_temp[2])
        elif len(txt_temp)==1:
            txt_box.append(txt_temp[0])
        else:
            txt_box.append('\n'.join(txt_temp))

    return txt_box
    #with open(os.path.join(txt_file_folder, txt_name),'w', encoding='utf-8') as f:
    #    f.write('\n'.join(txt_box))

def extract_song_from_ppt_file(ppt_file_path = "C:\\Users\\qiucanro\\Downloads\\2023-0723_checked.pptx", unique_keys=[]):
    #return a list of list item, each list item is of form [title, albulm, song_scripts]
    assert os.path.exists(ppt_file_path), "The specified ppt file path is not existing!"
    #_, name = os.path.split(ppt_file_path)
    #name = name.rsplit('.')[0]
    #assert os.path.exists(txt_file_folder), "the specified txt file folder is not existing!"
    ppt = Presentation(ppt_file_path)
    full_list = []
    slides = list(ppt.slides)
    song_list, title_pos = get_song_list(slides, unique_keys)
    slides = slides[title_pos+1:]
    for song_title in song_list:
        txt_box = [song_title]
        for i, slide in enumerate(slides):
            if len(slide.shapes)>=2 and hasattr(slide.shapes[0],'text') and hasattr(slide.shapes[1],'text') and \
                    (slide.shapes[0].text.strip().startswith(song_title) or slide.shapes[1].text.strip().startswith(song_title)):
                if slide.shapes[0].text.strip().startswith(song_title):
                    try:
                        txt_box.append(slide.shapes[1].text.strip())
                    except:
                        print('Something went wrong')
                elif slide.shapes[1].text.strip().startswith(song_title):
                    try:
                        txt_box.append(slide.shapes[1].text.strip().rsplit('\n')[1])
                    except:
                        print('Somthing went wrong!')
                if len(txt_box[-1].rsplit('\n'))>2:#very ugly solution to insert sign of no album
                    txt_box.insert(1,'Album not specified!')
                #now let's merge the first two items into one page section
                txt_box = ['\n'.join(txt_box)]
                move_on, j = True, 1
                while move_on:
                    if len(slides[i+j].shapes)>1:#and slides[i+j].shapes[0].text.startswith(song_title):
                        if len(list(slides[i+j].shapes)[0].text.strip())>len(list(slides[i+j].shapes)[1].text.strip()):
                            txt_box.append(list(slides[i+j].shapes)[0].text.strip().replace('\u3000','　').replace('\x0b','　'))
                        else:
                            txt_box.append(list(slides[i+j].shapes)[1].text.strip().replace('\u3000','　').replace('\x0b','　'))
                        j+=1
                    else:
                        move_on = False
                break
        full_list.append({'song_id':song_title, 'script':'\n\n'.join(txt_box)})
        unique_keys.append(song_title)
    return full_list, unique_keys

def get_song_list(slides, unique_keys=[]):
    songs = []
    title_pos = None
    for i, slide in enumerate(slides):
        if len(slide.shapes)>=2:
            if list(slide.shapes)[0].text.strip().startswith('诗歌赞美'):
                songs = songs + list(slide.shapes)[1].text.strip().rsplit('\n')
                title_pos = i
            elif list(slide.shapes)[0].text.strip().startswith('回应诗歌'):
                songs = songs + [list(slide.shapes)[1].text.strip().rsplit('\n')[0]]
    songs = [song for song in songs if song not in unique_keys]
    return songs, title_pos

def extract_songs_from_folder(folder = "C:\\Users\\qiucanro\\Downloads\\test", unique_keys = None):
    assert os.path.exists(folder), "The specified ppt file folder is not existing!"
    if unique_keys==None:
        unique_keys = []
    song_list = []
    # for each in glob.glob(os.path.join(folder, '*.pptx')):
    for each in glob.glob(os.path.join(folder, '**', '*.ppt'), recursive=True):
        try:
            songs, unique_keys = extract_song_from_ppt_file(each, unique_keys)
            song_list=song_list + songs
        except:
            print(f'Something went wrong in the process of extracting {each}, skip this file')
    with open(os.path.join(folder, 'songs.txt'),'w') as f:
        for each in song_list:
            f.write(each['script'])
            f.write('\n')
    return song_list

def save_ppt_content_to_json(folder="C:\\Users\\qiucanro\\Downloads\\Responsive Reading", json_file='scriptures.json'):
    assert os.path.exists(folder), "The specified ppt file folder is not existing!"
    results = {}
    for each in glob.glob(os.path.join(folder, '*.pptx')):
        temp = extract_text_from_ppt_file(each)
        results[temp[0].lstrip(' ').replace('\n','')] = '\n'.join(temp[1:])
    with open(os.path.join(folder, json_file),'w', encoding='utf-8') as f:
        json.dump(results, f)

def save_ppt_to_pptx_in_folder(folder):
    assert os.path.exists(folder), "The specified ppt file folder is not existing!"
    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True
    for each in glob.glob(os.path.join(folder, '*.ppt')):
        print(each)
        PPtPresentation = PptApp.Presentations.Open(each)
        PPtPresentation.SaveAs(each+'x', 24)
        PPtPresentation.close()
        #save_ppt_as_pptx(each)
    PptApp.Quit()

def test_extract_text_from_ppt_file():
    ppt_path = "C:\\Users\\qiucanro\\Downloads\\Responsive Reading\\02慈爱广大.pptx"
    txt_folder = "C:\\Users\\qiucanro\\Downloads"
    txt_name = '02慈爱广大.txt'
    extract_text_from_ppt_file(ppt_path, txt_folder, txt_name)
