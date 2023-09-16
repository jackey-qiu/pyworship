import os, json, re
import click
from pptx import Presentation
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

root = Path(__file__).parent.parent

class MakeWorkshipPpt(object):
    max_char_per_slide = 60
    use_json_for_extracting_scripture = True
    use_json_for_extracting_response_scripture = True
    font_type_kaiti = '方正楷体简体'#'FZKai-Z03S'
    font_type_zhunyuan = '方正准圆简体'#'FZZhunYuan-M02S'

    def __init__(self, date, holy_dinner_song_v = '2', content_folder = None):
        self.date = date
        self.holy_dinner_song_v = holy_dinner_song_v
        if len(date.rsplit('-')[-1])==2 and date.rsplit('-')[-1].startswith('0'):
            self.holy_dinner = int(date.rsplit('-')[-1][1:])<=7
        elif len(date.rsplit('-')[-1])==1:
            self.holy_dinner = int(date.rsplit('-')[-1])<=7
        else:
            self.holy_dinner = False
        self.src_folder = root / 'src' /'bkg_slides' / 'others'
        if content_folder==None:
            self.content_folder = root / 'src' / 'contents'
            self.ppt_file = root / 'src' / 'ppt_files' / f'{date}.pptx'
            self.content_folder.mkdir(parents=True, exist_ok=True)
            self.ppt_file.mkdir(parents=True, exist_ok=True)
        else:
            self.content_folder = content_folder
            Path(self.content_folder).mkdir(parents = True, exist_ok = True)
            self.ppt_file = Path(self.content_folder) / f'{date}.pptx'
        self.bible_json = root / 'src' / 'bible' / 'chinese_bible.json'
        self.scripture_json = root / 'src' / 'bible' / 'scriptures.json'
        self.prs = Presentation()
        self.prepare_slide_contents()

    def save_ppt_file(self):
        self.prs.save(self.ppt_file)

    def prepare_slide_contents(self):
        self._prepare_content_list('worker_list')
        self._prepare_content_list('report_list')       
        self._prepare_content_list('pray_list')
        self._prepare_content_list('song_list')
        if self.use_json_for_extracting_scripture:
            self._prepare_scripture_list_from_json()
        else:
            self._prepare_content_list('scripture_list')
        self._prepare_content_list('preach_list')

    def _prepare_content_list(self, attr_str):
        setattr(self, attr_str, [])
        yyyy, mm, dd = self.date.rsplit('-')
        file = os.path.join(self.content_folder, f'{attr_str}_{yyyy}-{mm}-{dd}.txt')
        if not os.path.exists(file):
            print(f'{file} is not existing!')
            return        
        with open(file, 'r', encoding= 'utf8') as f:
            lines = f.readlines()
            section_index_list = [i for i in range(len(lines)) if lines[i].startswith('#')]+[len(lines)]
            section_index_range = [(section_index_list[i], section_index_list[i+1]) for i in range(len(section_index_list)-1)]
            for i, ix_range in enumerate(section_index_range):
                lf, rt = ix_range
                getattr(self, attr_str).append([each.rstrip() for each in lines[lf+1:rt]])

    def _prepare_scripture_list_from_json(self, json_file = None, json_file_rep_scripture = None):
        #here you only provide book chapter verse signiture without having to give all scripture in the txt file
        #eg. book1:4[1-5,9,10]+book2:8[5,9,15-20]+book3:9[*]
        #extract vers 1 to 5 and ver 9 and ver 10 for chapter 4 in book1, 
        #extract vers 5, 9, and vers 15 to ver 20 for chapter 8 in book2,
        #extract all vers for chapter 9 in book3
        if json_file==None:
            json_file = self.bible_json
        if json_file_rep_scripture == None:
            json_file_scripture = self.scripture_json

        def _get_responsive_scripture_from_json(title):
            with open(json_file_scripture, 'r') as f:
                data = json.load(f)
                if title in data:
                    return [re.search(r'([0-9]?[0-9]?)(.*)',title).group(2)] + data[title].replace('\u3000','').replace('♂','').rsplit('\n')

        def _get_scripture_from_json(content_sig, json_obj):
            title_info = []
            
            scripture_list = []
            sections = content_sig.rstrip().rsplit('+')
            for section in sections:
                book, location = section.rsplit(':')
                chapter, vers_str = location.rsplit('[')
                ver_list_raw = vers_str[0:-1].rsplit(',')
                if vers_str[0:-1]=='*':
                    # title_info.append(f'{book}第{chapter}章')
                    title_info.append(f'{book}{chapter}')
                else:
                    # title_info.append(f'{book}第{chapter}章（{vers_str[0:-1]}）')
                    title_info.append(f'{book}{chapter}:{vers_str[0:-1]}')
                #vers = []
                for each in ver_list_raw:
                    if '-' not in each:
                        if each!='*':
                            scripture_list.append(each+json_obj[book][chapter][each].rstrip())
                        else:
                            #get full chapter scripture
                            scripture_list = [key+value.rstrip() for key, value in json_obj[book][chapter].items()]
                    else:
                        lf, rt = each.rsplit('-')
                        for i in range(int(lf),int(rt)+1):
                            scripture_list.append(str(i)+json_obj[book][chapter][str(i)].rstrip())
            return '；'.join(title_info), scripture_list

        #setattr(self, attr_str, [])
        self.scripture_list = []
        yyyy, mm, dd = self.date.rsplit('-')
        file = os.path.join(self.content_folder, f'scripture_list_{yyyy}-{mm}-{dd}.txt')
        if not os.path.exists(file):
            print(f'{file} is not existing!')
            return        
        with open(file, 'r', encoding= 'utf8') as f:
            lines = f.readlines()
            section_index_list = [i for i in range(len(lines)) if lines[i].startswith('#')]+[len(lines)]
            section_index_range = [(section_index_list[i], section_index_list[i+1]) for i in range(len(section_index_list)-1)]
            for i, ix_range in enumerate(section_index_range):
                lf, rt = ix_range
                self.scripture_list.append([each.rstrip() for each in lines[lf+1:rt]])

        with open(json_file, 'r') as f:
            bible = json.load(f)
            #modify first list (xuan zao scripture)
            sig = self.scripture_list[0][0]
            title_info, scripture_list = _get_scripture_from_json(sig, bible)
            self.scripture_list[0] = [title_info] + scripture_list
            #modify third list
            sig = self.scripture_list[2][0]
            title_info, scripture_list = _get_scripture_from_json(sig, bible)
            self.scripture_list[2] = [title_info] + scripture_list
            #modify second list
            second_list = _get_responsive_scripture_from_json(self.scripture_list[1][0])
            if second_list != None:
                self.scripture_list[1] = second_list

    def add_empty_slide(self, bkg_img = None):
        self.make_one_slide(blocks = [], bkg_img= bkg_img)

    def prepare_workship_slides(self):
        print('Prepare slides for pray section')
        self.prepare_pray_slides_all()
        self.add_empty_slide()
        print('Prepare slides for worker')
        self.prepare_worker_slide(which = 0)
        self.prepare_begin_slides()
        self.prepare_xuanzao_slides()
        self.add_empty_slide()
        print('Prepare slides for song section')
        self.prepare_song_slides(response=False)
        self.prepare_slides_for_response_scripture()
        self.prepare_main_pray()
        self.add_empty_slide()
        print('Prepare slides for scripture reading section')
        self.prepare_scripture_reading_slides()
        self.add_empty_slide()
        print('Prepare slides for preaching section')
        self.pepare_preach_slides()
        self.add_empty_slide()
        print('Prepare slides for response song section')
        self.prepare_song_slides(response=True)
        self.prepare_offering_slides()
        self.add_empty_slide()
        print('Prepare slides for report section')
        self.prepare_report_slides()
        self.add_empty_slide()
        self.prepare_new_friend_slide()
        self.prepare_worker_slide(which = 1)
        self.add_empty_slide()
        print('Prepare slides for finishing section')
        self.prepare_end_slides()
        self.save_ppt_file()

    def prepare_pray_slides_all(self):
        pray_string_list = self.pray_list
        for pray_string in pray_string_list:
            first_section = not bool(pray_string_list.index(pray_string))
            self.prepare_pray_slides_section(pray_string=pray_string, first_section=first_section)

    def prepare_worker_slide(self, which = 0):
        if which not in [0, 1]:
            print(f'There are two worker list maximum. But the index is {which}.')
            return
        content_format = {'cont':self.worker_list[which][1:],
                          'font_global':f'{self.font_type_kaiti}'+'+32+True',
                          'font_run':f'{self.font_type_kaiti}+24+True',
                          'alignment':'CENTER+37+7.68+0+0',
                          'textbox':['Cm(0.14)', 'Cm(1.16)', 'Cm(25.07)', 'Cm(17.27)']}
        title_format = {'cont':self.worker_list[which][0:1],
                          'font_global':f'{self.font_type_kaiti}'+'+24+True',
                          'font_run':f'{self.font_type_kaiti}'+'+24+True',
                          'alignment':'CENTER+28+0.25+0+0',
                          'textbox':['Cm(0.33)', 'Cm(0.09)', 'Cm(24.69)', 'Cm(1.45)']}
        self.make_one_slide(blocks = [content_format, title_format], bkg_img=os.path.join(self.src_folder, 'bkg_worker_list.jpg'))

    def prepare_begin_slides(self):
        title_format = {'cont':'',
                          'font_global':f'{self.font_type_kaiti}'+'+33+True',
                          'font_run':f'{self.font_type_kaiti}'+'+33+True',
                          'alignment':'LEFT+39+0.24+0+0',
                          'textbox': ['Cm(1.23)','Cm(5.6)','Cm(14.35)','Cm(1.8)']}
        title = '{} 年 {} 月 {} 日'.format(*self.date.rsplit('-'))
        title_format['cont'] = [title]        
        self.make_one_slide(blocks = [],bkg_img=os.path.join(self.src_folder, 'begin_slide_1.jpg'))
        self.make_one_slide(blocks = [title_format],bkg_img=os.path.join(self.src_folder, 'begin_slide_2.jpg'), middle_vertical=False)
        self.make_one_slide(blocks = [],bkg_img=os.path.join(self.src_folder, 'begin_slide_3.jpg'))
        self.make_one_slide(blocks = [],bkg_img=os.path.join(self.src_folder, 'begin_slide_4.jpg'))

    def prepare_xuanzao_slides(self):
        title = [self.scripture_list[0][0], '', '宣  召']
        self.prepare_slides_for_scripture(title, scriptures=self.scripture_list[0][1:], footnote_header='{}')

    def prepare_song_slides(self, response = False):
        #song_list = [[title, subtitle, scripts]]
        if not response:
            song_list = self.song_list[0:-1]
        else:
            song_list = [self.song_list[-1]]
        title_format = {'cont':[each[0] for each in song_list],
                        'font_global':'方正楷体简体+66+True',
                        'font_run':'方正楷体简体+66+True',
                        'alignment':'LEFT+77+0+0+0',
                        'textbox': ['Cm(1.13)','Cm(5.1)','Cm(23.6)','Cm(13.43)']}

        header_format = {'cont':[['诗歌赞美'],['回应诗歌']][int(response)],
                           'font_global':f'{self.font_type_kaiti}'+'+40+True',
                            'font_run':f'{self.font_type_kaiti}'+'+40+True',
                            'alignment':'LEFT+47+0+0+0',
                            'textbox': ['Cm(1.14)','Cm(0.51)','Cm(18.2)','Cm(1.91)']}

        subtitle_format = {'cont':'',
                            'font_global':'方正准圆简体+40+False',
                            'font_run':'方正准圆简体+40+False',
                            'alignment':'LEFT+47+0+0+0',
                            'textbox': ['Cm(1.33)','Cm(7.93)','Cm(20.07)','Cm(1.97)']}
         
        slide = self.make_one_slide(blocks = [title_format,header_format], bkg_img = os.path.join(self.src_folder,'bkg_general.jpg'), middle_vertical=False)
        if response:
            subtitle_format['cont'] = [song_list[0][1]]
            self.make_one_slide(slide = slide, blocks = [subtitle_format], middle_vertical=False)

        for song in song_list:
            tt, subtt = song[0:2]
            scripts = '\n'.join(song[2:])
            self.prepare_slides_for_one_song(title=[tt, subtt], scripts= scripts, include_title_page= not response)
            self.add_empty_slide()

    def prepare_main_pray(self):
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'main_pray_bkg_1.jpg'))
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'main_pray_bkg_2.jpg'))

    def prepare_new_friend_slide(self):
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'new_friends.jpg'))

    def prepare_offering_slides(self):
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'offering_1.jpg'))
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'offering_2.jpg'))
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'offering_3.jpg'))

    def prepare_report_slides(self):
        self.prepare_general_report_slides()
        content_format = {'cont':'',
                            'font_global':'方正楷体简体+44+True',
                            'font_run':'方正楷体简体+44+True',
                            'alignment':'LEFT+52+12+0+0',
                            'textbox': ['Cm(1.44)','Cm(1.44)','Cm(22.91)','Cm(15.72)']}

        content_footnote_format = {'cont':['活动报告'],
                            'font_global':'方正楷体简体+24+True',
                            'font_run':'方正楷体简体+24+True',
                            'alignment':'RIGHT+28+0+0+0',
                            'textbox': ['Cm(1.4)','Cm(17.5)','Cm(22.91)','Cm(1.1)']} 
        for report in self.report_list:
            content_format['cont'] = report  
            self.make_one_slide(blocks= [content_format, content_footnote_format], middle_vertical=False)           

    def prepare_general_report_slides(self):
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'report_1.jpg'))
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'report_2.jpg'))
        self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,'report_3.jpg'))

    def prepare_end_slides(self):
        if self.holy_dinner:
            if self.holy_dinner_song_v == '1':
                for i in range(1,23):
                    self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder.parent, 'holy_dinner_option1', f'end_{i}.jpg'))
            elif self.holy_dinner_song_v == '2':
                for i in range(1,21):
                    self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder.parent, 'holy_dinner_option2', f'Slide{i}.jpg'))
        else:
            for i in range(16,23):
                self.make_one_slide(blocks = [], bkg_img = os.path.join(self.src_folder,f'end_{i}.jpg'))

    def pepare_preach_slides(self):
        title, preacher, subtitle = self.preach_list[0]
        introduction, explanation, conclusion = self.preach_list[1:]
        title_format = {'cont':[title],
                            'font_global':f'{self.font_type_kaiti}'+'+66+True',
                            'font_run':f'{self.font_type_kaiti}'+'+66+True',
                            'alignment':'LEFT+77+0.24+0+0',
                            'textbox': ['Cm(1.14)','Cm(5.11)','Cm(22.55)','Cm(12.76)']}

        subtitle_format = {'cont':[preacher, subtitle],
                            'font_global':f'{self.font_type_kaiti}'+'+40+False',
                            'font_run':f'{self.font_type_kaiti}'+'+40+False',
                            'alignment':'LEFT+47+0+0+0',
                            'textbox':['Cm(1.14)','Cm(8.23)','Cm(22.55)','Cm(9.64)']}   

        content_format = {'cont':'',
                            'font_global':f'{self.font_type_kaiti}'+'+44+True',
                            'font_run':f'{self.font_type_kaiti}'+'+44+True',
                            'alignment':'LEFT+52+12+0+0',
                            'textbox': ['Cm(1.44)','Cm(0.73)','Cm(22.91)','Cm(16.43)']}

        content_footnote_format = {'cont':[f'{title}（{subtitle}）'],
                            'font_global':f'{self.font_type_kaiti}'+'+24+True',
                            'font_run':f'{self.font_type_kaiti}'+'+24+True',
                            'alignment':'RIGHT+28+0+0+0',
                            'textbox': ['Cm(1.4)','Cm(17.5)','Cm(22.91)','Cm(1.1)']}             
        
        max_chars_per_slide = self.max_char_per_slide
        def _prepare_slides_autofit_size(items, header = '', font_italic = False):
            previous_ix = 0
            for i in range(1, len(items)):
                overflow = sum([len(items[j]) for j in range(previous_ix, i)])>max_chars_per_slide 
                if overflow:
                    cont = items[previous_ix: i]
                    if header!='':
                        cont = [header] + cont
                    for ii in range(2,len(cont)+1):
                        content_format['cont'] = cont[0:ii]
                        slide = self.make_one_slide(blocks = [content_format], middle_vertical=False, font_italic=font_italic)
                        self.make_one_slide(slide = slide, blocks = [content_footnote_format], middle_vertical=False)
                    previous_ix = i
            #make one more slide at the end boundary
            cont = items[previous_ix: len(items)]
            if header!='':
                cont = [header] + cont
            for ii in range(2, len(cont)+1):
                content_format['cont'] = cont[0:ii]
                slide = self.make_one_slide(blocks = [content_format], middle_vertical=False, font_italic=font_italic)
                self.make_one_slide(slide = slide, blocks = [content_footnote_format], middle_vertical=False)

        self.make_one_slide(blocks= [title_format, subtitle_format], bkg_img=os.path.join(self.src_folder, 'preach_title_bkg.jpg'), middle_vertical=False)
        #introduction_items = introduction
        #explanation_items = explanation
        #conclusion_items = conclusion
        _prepare_slides_autofit_size(introduction, header = '引言：')
        _prepare_slides_autofit_size(explanation, header = '经文的理解与应用：')
        _prepare_slides_autofit_size(conclusion, header = '总结：', font_italic= True)

    def prepare_pray_slides_section(self, pray_string, first_section = True):
        title_format = {'cont':'',
                          'font_global':f'{self.font_type_kaiti}'+'+36+True',
                          'font_run':f'{self.font_type_kaiti}'+'+36+True',
                          'alignment':'LEFT+42+12+0+0',
                          'textbox': ['Cm(1.23)','Cm(5.6)','Cm(14.99)','Cm(1.8)']}
        
        content_format = {'cont':'',
                          'font_global':f'{self.font_type_kaiti}'+'+44+True',
                          'font_run':f'{self.font_type_kaiti}'+'+48+True',
                          'alignment':'LEFT+52+12+0+0',
                          'textbox': ['Cm(1.44)','Cm(1.44)','Cm(22.91)','Cm(15.72)']}
        
        footnote_format = {'cont':'',
                           'font_global':f'{self.font_type_kaiti}'+'+24+True',
                           'font_run':f'{self.font_type_kaiti}'+'+24+True',
                           'alignment':'RIGHT+28+0+0+0',
                           'textbox': ['Cm(1.4)','Cm(17.5)','Cm(22.91)','Cm(1.1)']}
        
        if first_section:
            title = '{} 年 {} 月 {} 日'.format(*self.date.rsplit('-'))
            title_format['cont'] = [title]
            self.make_one_slide(blocks = [], bkg_img= os.path.join(self.src_folder, 'bkg_pray_begin_slide.jpg'))
            self.make_one_slide(blocks = [title_format],bkg_img=os.path.join(self.src_folder, 'bkg_pray_title.jpg'), middle_vertical=False)
        footnote_format['cont'] = ['代祷事项']
        pray_items = pray_string
        previous_ix = 0
        if first_section:
            max_chars_per_slide = 1
        else:
            max_chars_per_slide = 50
        for i in range(1, len(pray_items)):
            if sum([len(pray_items[j]) for j in range(previous_ix, i)])>max_chars_per_slide:
                content_format['cont'] = pray_items[previous_ix: i]
                slide = self.make_one_slide(blocks = [content_format, footnote_format], middle_vertical=False)
                previous_ix = i
        #make one more slide at the end boundary
        cont = pray_items[previous_ix: len(pray_items)]
        content_format['cont'] = cont
        slide = self.make_one_slide(blocks = [content_format, footnote_format], middle_vertical=False)

    def prepare_slides_for_one_song(self, title, scripts, lines_per_page = 5, include_title_page = True):
        
        title_format = {'cont':'',
                            'font_global':f'{self.font_type_zhunyuan}'+'+80+True',
                            'font_run':f'{self.font_type_zhunyuan}'+'+80+True',
                            'alignment':'CENTER+94+0+0+0',
                            'textbox': ['Cm(1.75)','Cm(5.24)','Cm(21.91)','Cm(3.68)']}

        subtitle_format = {'cont':'',
                            'font_global':'方正准圆简体+40+False',
                            'font_run':'方正准圆简体+40+False',
                            'alignment':'CENTER+47+0+0+0',
                            'textbox': ['Cm(0.51)','Cm(8.93)','Cm(24.32)','Cm(1.97)']}

        script_format = {'cont':'',
                            'font_global':f'{self.font_type_zhunyuan}'+'+48+False',
                            'font_run':f'{self.font_type_zhunyuan}'+'+48+False',
                            'alignment':'CENTER+70+0+0+0',
                            'textbox': ['Cm(1.03)','Cm(1.13)','Cm(23.87)','Cm(15.6)']}

        script_footnote_format = {'cont':'',
                            'font_global':f'{self.font_type_kaiti}'+'+24+False',
                            'font_run':f'{self.font_type_kaiti}'+'+24+False',
                            'alignment':'RIGHT+36+0+0+0',
                            'textbox': ['Cm(0.93)','Cm(17.22)','Cm(24.06)','Cm(1.45)']}

        tt, subtt = title
        title_format['cont'] = [tt]
        subtitle_format['cont'] = [subtt]
        if include_title_page:
            self.make_one_slide(blocks = [title_format,subtitle_format])

        script_items = scripts.rsplit('\n')
        if '' in script_items:
            if script_items[0]!='':
                script_items = [''] + script_items
            if script_items[-1]!='':
                script_items.append('')
            section_ix_range = []
            section_ix_pair = []
            for i, each in enumerate(script_items):
                if each=='':
                    section_ix_range.append(i+1)
            #section_ix_range.append(len(script_items))
            for i in range(len(section_ix_range)-1):
                section_ix_pair.append([section_ix_range[i],section_ix_range[i+1]-1])
            pages = len(section_ix_pair)
            for i, pair in enumerate(section_ix_pair):
                script_format['cont'] = script_items[pair[0]:pair[1]]
                script_footnote_format['cont'] = [f'{tt}{i+1}/{pages}']
                self.make_one_slide(blocks = [script_format, script_footnote_format])
        else:
            pages = int((len(script_items) - len(script_items)%lines_per_page)/lines_per_page) + int((len(script_items)%lines_per_page)!=0)
            for i in range(0, len(script_items), lines_per_page):
                script_format['cont'] = script_items[i:(i+5)]
                script_footnote_format['cont'] = [f'{tt}{int(i/5+1)}/{pages}']
                self.make_one_slide(blocks = [script_format, script_footnote_format])

    def prepare_slides_for_response_scripture(self):
        title = self.scripture_list[1][0]
        scriptures = '\n'.join(self.scripture_list[1][1:])
        footnote_header = '启应经文（{}）'
        title_format = {'cont':'',
                        'font_global':'方正楷体简体+80+True',
                        'font_run':'方正楷体简体+80+True',
                        'alignment':'CENTER+94+0.24+0+0',
                        'textbox': ['Cm(0)','Cm(8.46)','Cm(25.4)','Cm(5.55)']}

        scripture_format = {'cont':'',
                            'font_global':'方正楷体简体+44+True',
                            'font_run':'方正楷体简体+44+True',
                            'alignment':'LEFT+52+12+0+0',
                            'textbox': ['Cm(1.44)','Cm(1.44)','Cm(22.91)','Cm(15.72)']}

        scripture_footnote_format = {'cont':'',
                            'font_global':'方正楷体简体+24+True',
                            'font_run':'方正楷体简体+24+True',
                            'alignment':'RIGHT+28+0+0+0',
                            'textbox': ['Cm(1.4)','Cm(17.5)','Cm(22.91)','Cm(1.1)']}    
            
        title_format['cont'] = [title]
        tt_footnote = title
        self.make_one_slide(blocks = [title_format],bkg_img=os.path.join(self.src_folder,'bkg_response_title.jpg'), middle_vertical=False, color_rgb=[255,255,0])
        scripture_footnote_format['cont'] = [footnote_header.format(tt_footnote)]
        scripture_items = scriptures.replace("神","　神").rsplit('\n')
        for i in range(0, int(len(scripture_items)/2) + int((len(scripture_items)%2)!=0)):
            scripture_on_slide = scripture_items[i*2:i*2+2]
            scripture_format['cont'] = scripture_on_slide
            slide = self.make_one_slide(blocks = [scripture_format], bkg_img=os.path.join(self.src_folder,'bkg_response_body.jpg'), middle_vertical=False, color_rgb=[[255,255,255],[255,255,0]])
            #if len(scripture_on_slide)==2:
            #    scripture_format['cont'] = [scripture_on_slide[1]]
            #    self.make_one_slide(slide = slide, blocks = [scripture_format], middle_vertical=False, color_rgb=[255,255,0])
            self.make_one_slide(slide = slide, blocks = [scripture_footnote_format], middle_vertical=False, color_rgb=[255,255,255])
        self.add_empty_slide(bkg_img=os.path.join(self.src_folder,'bkg_response_blank.jpg'))

    def prepare_scripture_reading_slides(self):
        title = [self.scripture_list[2][0], '', '读  经']
        scriptures = self.scripture_list[2][1:]
        self.prepare_slides_for_scripture(title, scriptures)

    def prepare_slides_for_scripture(self, title, scriptures, max_chars_per_slide = 60,footnote_header = '读经（{}）'):
        title_format = {'cont':'',
                        'font_global':'方正楷体简体+72+True',
                        'font_run':'方正楷体简体+72+True',
                        'alignment':'LEFT+83+0+0+0',
                        'textbox': ['Cm(1.13)','Cm(5.1)','Cm(23.6)','Cm(13.43)']}

        subtitle_format = {'cont':'',
                           'font_global':f'{self.font_type_kaiti}'+'+40+True',
                            'font_run':f'{self.font_type_kaiti}'+'+40+True',
                            'alignment':'LEFT+47+0+0+0',
                            'textbox': ['Cm(1.14)','Cm(0.51)','Cm(18.2)','Cm(1.91)']}

        scripture_format = {'cont':'',
                            'font_global':'方正楷体简体+44+True',
                            'font_run':'方正楷体简体+44+True',
                            'alignment':'LEFT+52+12+0+0',
                            'textbox': ['Cm(1.44)','Cm(1.44)','Cm(22.91)','Cm(15.72)']}

        scripture_footnote_format = {'cont':'',
                            'font_global':'方正楷体简体+24+True',
                            'font_run':'方正楷体简体+24+True',
                            'alignment':'RIGHT+28+0+0+0',
                            'textbox': ['Cm(1.4)','Cm(17.5)','Cm(22.91)','Cm(1.1)']}    
            
        tt, chapter, subtt = title
        title_format['cont'] = [tt, chapter]
        subtitle_format['cont'] = [subtt]
        tt_footnote = (tt+chapter).replace('+','')
        self.make_one_slide(blocks = [title_format,subtitle_format],bkg_img=os.path.join(self.src_folder,'bkg_general.jpg'), middle_vertical=False)
        scripture_footnote_format['cont'] = [footnote_header.format(tt_footnote)]
        scripture_items = scriptures
        previous_ix = 0
        for i in range(1, len(scripture_items)):
            if sum([len(scripture_items[j]) for j in range(previous_ix, i)])>max_chars_per_slide:
                scripture_format['cont'] = scripture_items[previous_ix: i]
                slide = self.make_one_slide(blocks = [scripture_format], middle_vertical=False, superscript_first_char=True)
                self.make_one_slide(slide = slide, blocks = [scripture_footnote_format], middle_vertical=False)
                previous_ix = i
        #make one more slide at the end boundary
        cont = scripture_items[previous_ix: len(scripture_items)]
        scripture_format['cont'] = cont
        slide = self.make_one_slide(blocks = [scripture_format], middle_vertical=False, superscript_first_char=True)
        self.make_one_slide(slide = slide, blocks = [scripture_footnote_format], middle_vertical=False)

    def make_one_slide(self, blocks, slide = None, bkg_img = None, middle_vertical = True,superscript_first_char = False, color_rgb = None, font_italic = False):
        prs = self.prs
        color_index = -1
        color_circle = False
        if color_rgb!=None and type(color_rgb[0])==list:
            color_circle = True
        if slide == None:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        if bkg_img != None:
            slide.shapes.add_picture(bkg_img, 0, 0, height = prs.slide_height)
        for block in blocks:
            txBox = slide.shapes.add_textbox(*[eval(each) for each in block['textbox']])
            tf = txBox.text_frame
            text_pg = block['cont']
            tf.word_wrap = True
            if middle_vertical:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            for i,txt in enumerate(text_pg):
                if color_circle:
                    color_index = color_index + 1
                    if color_index==len(color_rgb):
                        color_index = 0
                if i==0:
                    p = tf.paragraphs[0]
                    print('processing paragraph 1')
                else:
                    p = tf.add_paragraph()
                    print(f'processing paragraph {i+1}')
                p.alignment = getattr(PP_PARAGRAPH_ALIGNMENT, block['alignment'].rsplit('+')[0])
                p.space_before = Pt(float(block['alignment'].rsplit('+')[2]))
                p.space_after = Pt(float(block['alignment'].rsplit('+')[3]))
                p.line_spacing = Pt(float(block['alignment'].rsplit('+')[1]))
                p.level = eval(block['alignment'].rsplit('+')[4])
                if superscript_first_char:
                    last_num_ix = 0
                    while True:
                        if txt[last_num_ix] not in [str(ii) for ii in range(10)]:
                            # last_num_ix = last_num_ix - 1
                            break
                        else:
                            last_num_ix = last_num_ix + 1
                    if last_num_ix == 0:
                        txt = txt[0]+'+'+txt[1:]
                    else:
                        txt = txt[0:last_num_ix]+'+'+txt[last_num_ix:]
                txt_runs = txt.rsplit('+')
                for j,txt_run in enumerate(txt_runs):
                    run = p.add_run()
                    run.text = txt_run
                    if j==0:
                        nm, sz, bold =  block['font_global'].rsplit('+')
                        if superscript_first_char:
                            run.font._element.set('baseline', '30000')
                    else:
                        nm, sz, bold =  block['font_run'].rsplit('+')
                    run.font.language_id = MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE
                    run.font.size = Pt(eval(sz))
                    run.font.name = nm
                    run.font.bold = eval(bold)
                    if color_rgb!=None:
                        if not color_circle:
                            run.font.color.rgb = RGBColor(*color_rgb)
                        else:
                            run.font.color.rgb = RGBColor(*color_rgb[color_index])
                    if font_italic:
                        run.font.italic = True
                # set language id 
                p.font.language_id = MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE
                # set font typeface for asian characters
                defrpr = p._element.pPr.defRPr
                ea = etree.SubElement(defrpr, qn('a:ea'))
                ea.set('typeface', nm)                        
        return slide

def main(date, v, content_folder):
    ppt_worker = MakeWorkshipPpt(date, v, content_folder=content_folder)
    ppt_worker.prepare_workship_slides()     

@click.command()
@click.argument('date', required=1)
@click.option('--v', default='2', help='song type: either 1 or 2')
def ppt(date, v):
    main(date, v)

if __name__ == '__main__':
    import sys
    song_v = '2'
    if len(sys.argv) == 2:
        date = sys.argv[1]
    elif len(sys.argv) == 3:
        date, song_v= sys.argv[1:]
    else:
        print("请输入主日崇拜日期，格式为yyyy-mm-dd，比如2023-08-16:")
        date = input()
        print("如果是圣餐主日，请输入圣餐赞美诗歌\n1(靠近十架) or 2 (宝架清影), 输入1或者2:")
        song_v = input()
    missing_files = []
    for each in ['pray_list_', 'preach_list_', 'song_list_', 'scripture_list_','report_list_','worker_list_']:
        if not os.path.exists(root / 'src' / 'contents' / f'{each}{date}.txt'):
            missing_files.append(root / 'src' / 'contents' / f'{each}{date}.txt')
    if len(missing_files)!=0:
        print('以下文件缺失，请加入后再运行程序。')
        for each in missing_files:
            print(each)
    else:
        ppt(date, song_v)